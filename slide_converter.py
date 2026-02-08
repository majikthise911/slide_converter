"""
convert — PDF/PPTX to structured HTML or Markdown with embedded images.

Uses font analysis for heading/body/math detection, bullet structure,
and equation identification. Auto-renders pages with diagrams or equations
as images to guarantee visual accuracy.

Usage:
    convert lecture.pdf                     # → lecture.html (default)
    convert lecture.pdf --md                # → lecture.md
    convert lecture.pdf -o notes.md         # auto-detects format from extension
    convert lecture.pdf --render            # render ALL pages as images
    convert lecture.pdf --no-render         # text only (smallest file)
    convert file1.pdf file2.pdf -o all.html # merge multiple files
"""
import sys
import os
import base64
import html as html_mod
from pathlib import Path
from collections import Counter

try:
    import pymupdf
except ImportError:
    print("Installing pymupdf...")
    os.system(f"{sys.executable} -m pip install -q pymupdf")
    import pymupdf

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# --- Constants ---------------------------------------------------------------

MATH_FONTS = ("CambriaMath", "Cambria Math", "MT-Extra", "Symbol")

RENDER_AUTO = "auto"
RENDER_ALL = "all"
RENDER_NONE = "none"

FMT_HTML = "html"
FMT_MD = "md"

# Element types
TITLE = "title"
BULLET = "bullet"
EQUATION = "equation"
CODE = "code"
LABEL = "label"
BODY = "body"
IMAGE = "image"
RENDER = "render"
TABLE = "table"

CSS = """\
body { font-family: -apple-system, 'Segoe UI', Arial, sans-serif; max-width: 960px; margin: 40px auto; padding: 20px; line-height: 1.7; color: #1a1a1a; }
h1 { border-bottom: 2px solid #333; padding-bottom: 10px; margin-bottom: 10px; }
nav { margin-bottom: 30px; padding: 15px; background: #f8f9fa; border-radius: 6px; }
nav summary { font-weight: bold; cursor: pointer; }
nav ol { columns: 2; column-gap: 2em; margin: 10px 0 0 0; padding-left: 1.5em; }
nav li { margin: 2px 0; font-size: 0.92em; break-inside: avoid; }
nav a { color: #2471a3; text-decoration: none; }
nav a:hover { text-decoration: underline; }
.slide { margin-bottom: 2em; }
.slide-title { color: #1a5276; margin-top: 2em; padding: 4px 0; border-bottom: 1px solid #ddd; }
ul { margin: 0.3em 0; padding-left: 1.8em; }
li { margin: 0.2em 0; line-height: 1.6; }
.math { font-family: 'Cambria Math', 'STIX Two Math', 'Latin Modern Math', serif; }
.eq { display: block; margin: 0.6em 1.5em; padding: 6px 12px; background: #f0f4f8; border-left: 3px solid #2980b9; font-family: 'Cambria Math', serif; font-size: 1.05em; white-space: pre-wrap; }
img { max-width: 100%; height: auto; margin: 1em 0; display: block; }
.slide-img { max-width: 100%; border: 1px solid #ddd; border-radius: 4px; margin: 0.5em 0; }
.label { font-size: 0.88em; color: #555; margin: 0.3em 0; }
pre, code { background: #f5f5f5; font-family: Menlo, Consolas, monospace; font-size: 0.9em; }
pre { padding: 12px; border-radius: 4px; overflow-x: auto; }
code { padding: 2px 5px; border-radius: 3px; }
table { border-collapse: collapse; margin: 1em 0; width: 100%; }
th, td { border: 1px solid #ccc; padding: 6px 10px; text-align: left; }
th { background: #f0f0f0; font-weight: bold; }
details.render { margin: 0.5em 0; }
details.render summary { cursor: pointer; color: #888; font-size: 0.82em; }
"""


# --- Helpers -----------------------------------------------------------------

def esc(text):
    return html_mod.escape(text)


def is_math(font_name):
    return any(m in font_name for m in MATH_FONTS)


def el(etype, **kw):
    """Create an element dict."""
    return {"type": etype, **kw}


def plain_text(spans):
    """Extract plain text from spans."""
    return "".join(s["text"] for s in spans).strip()


def strip_bullet_char(spans):
    """Return spans with the leading bullet character removed."""
    out = []
    stripped = False
    for s in spans:
        if not stripped:
            t = s["text"].lstrip()
            if t and t[0] in "\u2022\u2023\u25cf\u25cb\u2013\u2014\u2012":
                t = t[1:].lstrip()
                if t:
                    out.append({**s, "text": t})
                stripped = True
                continue
            elif t:
                out.append(s)
                stripped = True
                continue
        else:
            out.append(s)
    return out if out else spans


# --- Span renderers ---------------------------------------------------------

def spans_to_html(spans):
    """Render spans to HTML."""
    parts = []
    for s in spans:
        t = s["text"]
        if not t:
            continue
        e = esc(t)
        font = s["font"]
        if is_math(font):
            parts.append(f'<span class="math">{e}</span>')
        elif "Bold" in font and "Italic" in font:
            parts.append(f"<strong><em>{e}</em></strong>")
        elif "Bold" in font:
            parts.append(f"<strong>{e}</strong>")
        elif "Italic" in font:
            parts.append(f"<em>{e}</em>")
        else:
            parts.append(e)
    return "".join(parts)


def spans_to_md(spans):
    """Render spans to Markdown."""
    parts = []
    for s in spans:
        t = s["text"]
        if not t:
            continue
        font = s["font"]
        # Escape markdown special chars in text (but not formatting we add)
        escaped = t.replace("\\", "\\\\")
        for ch in "*_`[]()#>":
            escaped = escaped.replace(ch, f"\\{ch}")
        if is_math(font):
            parts.append(t)  # keep math Unicode as-is, no escaping
        elif "Bold" in font and "Italic" in font:
            parts.append(f"***{escaped}***")
        elif "Bold" in font:
            parts.append(f"**{escaped}**")
        elif "Italic" in font:
            parts.append(f"*{escaped}*")
        else:
            parts.append(escaped)
    return "".join(parts)


# --- PDF extraction (format-agnostic) ---------------------------------------

def analyze_pdf_fonts(doc):
    """Determine adaptive font-size thresholds."""
    size_chars = Counter()
    for page in doc:
        for block in page.get_text("dict")["blocks"]:
            if block["type"] != 0:
                continue
            for line in block["lines"]:
                for span in line["spans"]:
                    txt = span["text"].strip()
                    if txt:
                        size_chars[round(span["size"])] += len(txt)
    if not size_chars:
        return {"title": 40, "body": 24}
    sorted_sizes = sorted(size_chars.keys(), reverse=True)
    return {"title": sorted_sizes[0], "body": size_chars.most_common(1)[0][0]}


def page_needs_render(page):
    """Check if a page has diagrams or math needing a visual render."""
    if len(page.get_drawings()) > 4:
        return True
    for block in page.get_text("dict")["blocks"]:
        if block["type"] != 0:
            continue
        for line in block["lines"]:
            for span in line["spans"]:
                if span["text"].strip() and is_math(span["font"]):
                    return True
    return False


def postprocess_elements(elements):
    """Merge consecutive equations; detect code blocks from labels."""
    result = []
    i = 0
    while i < len(elements):
        e = elements[i]

        # Merge consecutive equations
        if e["type"] == EQUATION:
            merged_spans = list(e["spans"])
            j = i + 1
            while j < len(elements) and elements[j]["type"] == EQUATION:
                merged_spans.extend(elements[j]["spans"])
                j += 1
            if j > i + 1:
                result.append(el(EQUATION, spans=merged_spans))
            else:
                result.append(e)
            i = j
            continue

        # Detect code in consecutive labels
        if e["type"] == LABEL:
            j = i
            labels = []
            while j < len(elements) and elements[j]["type"] == LABEL:
                labels.append(elements[j])
                j += 1
            texts = [plain_text(lb["spans"]) for lb in labels]
            code_hits = sum(1 for t in texts if ";" in t or "=" in t or "(" in t)
            if len(labels) >= 3 and code_hits >= 2:
                result.append(el(CODE, lines=texts))
            else:
                result.extend(labels)
            i = j
            continue

        result.append(e)
        i += 1
    return result


def extract_pdf_page(page, page_num, doc, fs):
    """Extract a PDF page into a list of elements and a slide title."""
    elements = []
    d = page.get_text("dict")
    page_h = d["height"]

    list_depth = 0
    current_bullet = None  # index into elements of the open bullet
    title_emitted = False
    slide_title = None

    def close_bullet():
        nonlocal current_bullet
        current_bullet = None

    def close_lists():
        nonlocal list_depth
        close_bullet()
        list_depth = 0

    for block in d["blocks"]:
        if block["type"] != 0:
            continue
        for line in block["lines"]:
            spans = [s for s in line["spans"] if s["text"].strip()]
            if not spans:
                continue
            text = "".join(s["text"] for s in spans).strip()
            if not text:
                continue

            max_sz = max(s["size"] for s in spans)
            y = line["bbox"][1]
            all_math = all(is_math(s["font"]) for s in spans)

            # Skip page numbers
            if max_sz < 15 and y > page_h * 0.80 and text.strip().isdigit():
                continue

            # Title
            if max_sz >= fs["title"] - 2 and not title_emitted:
                close_lists()
                slide_title = text
                elements.append(el(TITLE, spans=list(spans), page_num=page_num))
                title_emitted = True
                continue

            # Top-level bullet
            if text[0] in "\u2022\u2023\u25cf\u25cb":
                close_bullet()
                if list_depth > 1:
                    list_depth = 1
                if list_depth == 0:
                    list_depth = 1
                current_bullet = len(elements)
                elements.append(el(BULLET, spans=strip_bullet_char(spans), level=0))
                continue

            # Sub-bullet
            if text[0] in "\u2013\u2014\u2012":
                close_bullet()
                if list_depth < 1:
                    list_depth = 1
                list_depth = 2
                current_bullet = len(elements)
                elements.append(el(BULLET, spans=strip_bullet_char(spans), level=1))
                continue

            # Equation
            if all_math:
                close_bullet()
                elements.append(el(EQUATION, spans=list(spans)))
                continue

            # Small text / label
            if max_sz < fs["body"] - 6 and max_sz < fs["title"] - 10:
                close_lists()
                elements.append(el(LABEL, spans=list(spans)))
                continue

            # Body / continuation
            if current_bullet is not None:
                elements[current_bullet]["spans"].extend(spans)
            elif list_depth > 0:
                current_bullet = len(elements)
                elements.append(el(BULLET, spans=list(spans), level=list_depth - 1))
            else:
                elements.append(el(BODY, spans=list(spans)))

    close_lists()

    # Raster images
    for idx, img_info in enumerate(page.get_images()):
        try:
            xref = img_info[0]
            bi = doc.extract_image(xref)
            elements.append(el(
                IMAGE,
                data=base64.b64encode(bi["image"]).decode(),
                ext=bi["ext"],
                alt=f"Slide {page_num + 1} Figure {idx + 1}",
            ))
        except Exception:
            pass

    elements = postprocess_elements(elements)
    return elements, slide_title


def extract_pdf(pdf_path, render_mode):
    """Extract all pages from a PDF. Returns (toc, list_of_page_elements)."""
    doc = pymupdf.open(str(pdf_path))
    fs = analyze_pdf_fonts(doc)
    n = len(doc)
    print(f"  {n} pages | title={fs['title']}pt body={fs['body']}pt")

    toc = []
    pages = []
    rendered = 0

    for pn in range(n):
        page = doc[pn]
        elems, title = extract_pdf_page(page, pn, doc, fs)

        should_render = (
            render_mode == RENDER_ALL
            or (render_mode == RENDER_AUTO and page_needs_render(page))
        )
        if should_render:
            pix = page.get_pixmap(dpi=120)
            elems.append(el(
                RENDER,
                data=base64.b64encode(pix.tobytes("png")).decode(),
                alt=f"Slide {pn + 1}",
            ))
            rendered += 1

        toc.append((pn + 1, title or f"Slide {pn + 1}"))
        pages.append(elems)
        print(f"  Page {pn + 1}/{n}\r", end="", flush=True)

    doc.close()
    print()
    if rendered:
        print(f"  {rendered}/{n} pages rendered as images")
    return toc, pages


# --- PPTX extraction (format-agnostic) --------------------------------------

def extract_pptx(pptx_path):
    """Extract all slides from a PPTX. Returns (toc, list_of_slide_elements)."""
    if not HAS_PPTX:
        print("  ERROR: python-pptx not installed. Run: pip install python-pptx")
        return [], []

    prs = Presentation(str(pptx_path))
    n = len(prs.slides)
    print(f"  {n} slides")

    toc = []
    pages = []

    for sn, slide in enumerate(prs.slides, 1):
        elems = []
        slide_title = None
        shapes = sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0))

        for shape in shapes:
            # Tables
            if shape.has_table:
                tbl = shape.table
                rows = []
                for row in tbl.rows:
                    rows.append([cell.text for cell in row.cells])
                elems.append(el(TABLE, rows=rows))
                continue

            # Images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    elems.append(el(
                        IMAGE,
                        data=base64.b64encode(shape.image.blob).decode(),
                        ext=shape.image.ext,
                        alt=f"Slide {sn} Image",
                    ))
                except Exception:
                    pass
                continue

            # Text
            if not shape.has_text_frame or not shape.text.strip():
                continue

            is_title_shape = False
            try:
                if shape.placeholder_format is not None:
                    is_title_shape = shape.placeholder_format.idx in (0, 1)
            except Exception:
                pass

            for para in shape.text_frame.paragraphs:
                if not para.text.strip():
                    continue

                # Build fake spans from runs (same structure as PDF spans)
                spans = []
                for run in para.runs:
                    if not run.text:
                        continue
                    font_name = "Normal"
                    try:
                        if run.font.bold and run.font.italic:
                            font_name = "BoldItalic"
                        elif run.font.bold:
                            font_name = "Bold"
                        elif run.font.italic:
                            font_name = "Italic"
                    except Exception:
                        pass
                    spans.append({"text": run.text, "font": font_name, "size": 24})

                if not spans:
                    spans = [{"text": para.text, "font": "Normal", "size": 24}]

                if is_title_shape and slide_title is None:
                    slide_title = para.text.strip()
                    elems.append(el(TITLE, spans=spans, page_num=sn - 1))
                    is_title_shape = False
                elif para.level > 0:
                    elems.append(el(BULLET, spans=spans, level=min(para.level, 1)))
                else:
                    elems.append(el(BODY, spans=spans))

        toc.append((sn, slide_title or f"Slide {sn}"))
        pages.append(elems)
        print(f"  Slide {sn}/{n}\r", end="", flush=True)

    print()
    return toc, pages


# --- HTML renderer -----------------------------------------------------------

def elements_to_html(elements):
    """Render a list of elements to HTML."""
    parts = []
    in_list = 0  # nesting depth

    def ensure_list(level):
        nonlocal in_list
        target = level + 1
        while in_list < target:
            parts.append("<ul>")
            in_list += 1
        while in_list > target:
            parts.append("</ul>")
            in_list -= 1

    def close_list():
        nonlocal in_list
        while in_list > 0:
            parts.append("</ul>")
            in_list -= 1

    for e in elements:
        t = e["type"]

        if t == TITLE:
            close_list()
            pn = e["page_num"] + 1
            parts.append(
                f'<h2 class="slide-title" id="slide-{pn}">'
                f"{spans_to_html(e['spans'])}</h2>"
            )

        elif t == BULLET:
            ensure_list(e["level"])
            parts.append(f"<li>{spans_to_html(e['spans'])}</li>")

        elif t == EQUATION:
            close_list()
            parts.append(f'<div class="eq">{spans_to_html(e["spans"])}</div>')

        elif t == CODE:
            close_list()
            parts.append(
                "<pre><code>"
                + "\n".join(esc(line) for line in e["lines"])
                + "</code></pre>"
            )

        elif t == LABEL:
            close_list()
            parts.append(f'<p class="label">{spans_to_html(e["spans"])}</p>')

        elif t == BODY:
            close_list()
            parts.append(f"<p>{spans_to_html(e['spans'])}</p>")

        elif t == IMAGE:
            close_list()
            parts.append(
                f'<img src="data:image/{e["ext"]};base64,{e["data"]}" '
                f'alt="{esc(e["alt"])}">'
            )

        elif t == RENDER:
            close_list()
            parts.append(
                f'<img class="slide-img" src="data:image/png;base64,{e["data"]}" '
                f'alt="{esc(e["alt"])}">'
            )

        elif t == TABLE:
            close_list()
            parts.append("<table>")
            for ri, row in enumerate(e["rows"]):
                parts.append("<tr>")
                tag = "th" if ri == 0 else "td"
                for cell in row:
                    parts.append(f"<{tag}>{esc(cell)}</{tag}>")
                parts.append("</tr>")
            parts.append("</table>")

    close_list()
    return "\n".join(parts)


def assemble_html(title, toc, pages_html):
    """Build a full HTML document."""
    toc_html = '<nav><details open><summary>Table of Contents</summary><ol>'
    for pn, t in toc:
        toc_html += f'<li><a href="#slide-{pn}">{esc(t)}</a></li>'
    toc_html += "</ol></details></nav>"

    body = toc_html + "\n"
    for page_html in pages_html:
        body += f'<section class="slide">\n{page_html}\n</section>\n'

    return (
        f"<!DOCTYPE html>\n<html lang=\"en\">\n<head>\n"
        f"<meta charset=\"UTF-8\">\n<title>{esc(title)}</title>\n"
        f"<style>\n{CSS}</style>\n</head>\n<body>\n"
        f"<h1>{esc(title)}</h1>\n{body}</body>\n</html>"
    )


# --- Markdown renderer -------------------------------------------------------

def elements_to_md(elements):
    """Render a list of elements to Markdown."""
    lines = []

    for e in elements:
        t = e["type"]

        if t == TITLE:
            lines.append(f"## {spans_to_md(e['spans'])}")
            lines.append("")

        elif t == BULLET:
            indent = "  " * e["level"]
            lines.append(f"{indent}- {spans_to_md(e['spans'])}")

        elif t == EQUATION:
            lines.append("")
            lines.append(f"> {plain_text(e['spans'])}")
            lines.append("")

        elif t == CODE:
            lines.append("")
            lines.append("```")
            lines.extend(e["lines"])
            lines.append("```")
            lines.append("")

        elif t == LABEL:
            lines.append(f"*{spans_to_md(e['spans'])}*")

        elif t == BODY:
            lines.append("")
            lines.append(spans_to_md(e["spans"]))
            lines.append("")

        elif t == IMAGE:
            lines.append("")
            lines.append(f'![{e["alt"]}](data:image/{e["ext"]};base64,{e["data"]})')
            lines.append("")

        elif t == RENDER:
            lines.append("")
            lines.append(f'![{e["alt"]}](data:image/png;base64,{e["data"]})')
            lines.append("")

        elif t == TABLE:
            lines.append("")
            header = e["rows"][0] if e["rows"] else []
            lines.append("| " + " | ".join(header) + " |")
            lines.append("| " + " | ".join("---" for _ in header) + " |")
            for row in e["rows"][1:]:
                lines.append("| " + " | ".join(row) + " |")
            lines.append("")

    return "\n".join(lines)


def assemble_md(title, toc, pages_md):
    """Build a full Markdown document."""
    lines = [f"# {title}", ""]

    # Table of contents
    lines.append("## Table of Contents")
    lines.append("")
    for pn, t in toc:
        lines.append(f"{pn}. [{t}](#slide-{pn})")
    lines.append("")
    lines.append("---")
    lines.append("")

    for page_md in pages_md:
        lines.append(page_md)
        lines.append("")

    return "\n".join(lines)


# --- Main pipeline -----------------------------------------------------------

def convert_file(path, render_mode, fmt):
    """Convert a single file. Returns (content_string, filename)."""
    p = Path(path)
    ext = p.suffix.lower()
    print(f"Converting: {p.name}")

    if ext == ".pdf":
        toc, pages_elements = extract_pdf(p, render_mode)
    elif ext == ".pptx":
        toc, pages_elements = extract_pptx(p)
    else:
        print(f"  ERROR: Unsupported format '{ext}'. Use .pdf or .pptx")
        sys.exit(1)

    if fmt == FMT_HTML:
        pages_rendered = [elements_to_html(elems) for elems in pages_elements]
        return assemble_html(p.stem, toc, pages_rendered), p.stem
    else:
        pages_rendered = [elements_to_md(elems) for elems in pages_elements]
        return assemble_md(p.stem, toc, pages_rendered), p.stem


def main():
    args = sys.argv[1:]

    if not args or args[0] in ("-h", "--help"):
        print("Usage: convert [options] <file1> [file2 ...] [-o output.html|.md]")
        print()
        print("Options:")
        print("  --md          Output Markdown instead of HTML")
        print("  --render      Render ALL pages as images (largest file)")
        print("  --no-render   Text extraction only, no page renders (smallest file)")
        print("  -o FILE       Output file (.html or .md auto-detected from extension)")
        print()
        print("Default: HTML with auto-rendered pages for diagrams/equations")
        print()
        print("Examples:")
        print('  convert "SPCE5025 Class 3.pdf"')
        print('  convert lecture.pdf --md')
        print('  convert lecture.pdf --render')
        print('  convert lecture.pdf -o notes.md')
        print('  convert week1.pdf week2.pdf -o combined.html')
        sys.exit(0)

    # Parse flags
    use_md = "--md" in args
    args = [a for a in args if a != "--md"]

    if "--render" in args:
        render_mode = RENDER_ALL
        args = [a for a in args if a != "--render"]
    elif "--no-render" in args:
        render_mode = RENDER_NONE
        args = [a for a in args if a != "--no-render"]
    else:
        render_mode = RENDER_AUTO

    output_path = None
    if "-o" in args:
        oi = args.index("-o")
        if oi + 1 < len(args):
            output_path = args[oi + 1]
            args = args[:oi] + args[oi + 2:]
        else:
            print("Error: -o requires a filename")
            sys.exit(1)

    input_files = [Path(a) for a in args]
    for f in input_files:
        if not f.exists():
            print(f"Error: File not found: {f}")
            sys.exit(1)
    if not input_files:
        print("Error: No input files specified")
        sys.exit(1)

    # Determine output format
    if output_path and Path(output_path).suffix.lower() == ".md":
        fmt = FMT_MD
    elif use_md:
        fmt = FMT_MD
    else:
        fmt = FMT_HTML

    # Determine output path
    if output_path:
        out = Path(output_path)
    elif len(input_files) == 1:
        out = input_files[0].with_suffix(".md" if fmt == FMT_MD else ".html")
    else:
        out = Path("combined.md" if fmt == FMT_MD else "combined.html")

    # Convert
    sections = []
    for f in input_files:
        content, name = convert_file(f, render_mode, fmt)
        sections.append((content, name))

    # Merge if multiple files
    if len(sections) == 1:
        output = sections[0][0]
    else:
        if fmt == FMT_HTML:
            title = "Combined: " + ", ".join(s[1] for s in sections)
            merged = "\n<hr>\n".join(s[0] for s in sections)
            output = (
                f"<!DOCTYPE html>\n<html lang=\"en\">\n<head>\n"
                f"<meta charset=\"UTF-8\">\n<title>{esc(title)}</title>\n"
                f"<style>\n{CSS}</style>\n</head>\n<body>\n"
                f"<h1>{esc(title)}</h1>\n{merged}\n</body>\n</html>"
            )
        else:
            parts = []
            for content, name in sections:
                parts.append(f"# {name}\n\n{content}")
            output = "\n\n---\n\n".join(parts)

    with open(out, "w", encoding="utf-8") as f:
        f.write(output)

    size_mb = out.stat().st_size / 1024 / 1024
    print(f"\nDone: {out} ({size_mb:.1f} MB)")


if __name__ == "__main__":
    main()
